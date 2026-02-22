"""
╔═══════════════════════════════════════════════╗
║      🎓  AkademikYordamchi Bot  v3.0         ║
║   AI yordamida professional prezentatsiya    ║
║   Gemini 2.5 Flash  |  python-pptx           ║
║   📄 PDF / PPTX / DOCX → Slayd yangi!        ║
╚═══════════════════════════════════════════════╝
"""

import os, re, json, logging, asyncio, datetime, requests, tempfile
from io import BytesIO
from pathlib import Path
from functools import wraps

from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup, BotCommand
from telegram.ext import (
    ApplicationBuilder, CommandHandler, MessageHandler,
    CallbackQueryHandler, filters, ContextTypes
)
from telegram.constants import ParseMode, ChatAction

from google import genai
from dotenv import load_dotenv

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

try:
    from PIL import Image as PILImage
    PIL_OK = True
except ImportError:
    PIL_OK = False

# PDF matnini o'qish uchun
try:
    import fitz  # PyMuPDF
    PYMUPDF_OK = True
except ImportError:
    PYMUPDF_OK = False

# DOCX matnini o'qish uchun
try:
    from docx import Document as DocxDocument
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

# ════════════════════════════════════════════════
# LOGGING
# ════════════════════════════════════════════════
logging.basicConfig(
    format="%(asctime)s | %(levelname)s | %(name)s | %(message)s",
    level=logging.INFO
)
log = logging.getLogger(__name__)

# ════════════════════════════════════════════════
# ENV
# ════════════════════════════════════════════════
load_dotenv()
TOKEN      = os.getenv("TELEGRAM_TOKEN", "8566361305:AAGZARrVlfZoVATqez_MRTWehiZjWckVHGg")
GEMINI_KEY = os.getenv("GEMINI_API_KEY", "AIzaSyAEO8E7RKx_c76WFNE8mTNjVMW5X46xh3g")
PEXELS_KEY = os.getenv("PEXELS_API_KEY", "563492ad6f917000010000014744f333f2e647ffaa4c2609ec5be16f")
ADMIN_IDS  = {int(x) for x in os.getenv("ADMIN_IDS", "").split(",") if x.strip().isdigit()}

if not TOKEN or not GEMINI_KEY:
    raise ValueError(".env faylida TELEGRAM_TOKEN yoki GEMINI_API_KEY yoq!")

ai    = genai.Client(api_key=GEMINI_KEY)
MODEL = "gemini-2.5-flash"

# ════════════════════════════════════════════════
# LIMITLAR
# ════════════════════════════════════════════════
FREE_DAY   = 2
PREM_DAY   = 20
FREE_SLIDE = 8
PREM_SLIDE = 20
PRICE      = "19 900 so'm/oy"

S_IDLE       = "idle"
S_TOPIC      = "topic"
S_FILE_TOPIC = "file_topic"   # Fayl yuklangandan keyin mavzu so'rash

# ════════════════════════════════════════════════
# DIZAYN TEMALARI
# ════════════════════════════════════════════════
THEMES = {
    "classic": {"tbg":"1E3A5F","tfg":"FFFFFF","bbg":"F0F4FA","bfg":"1A1A2E","acc":"E63946","name":"🎓 Klassik Ko'k"},
    "dark":    {"tbg":"16213E","tfg":"E2E2E2","bbg":"0F3460","bfg":"E2E2E2","acc":"E94560","name":"🌙 Dark Pro"},
    "green":   {"tbg":"1B5E20","tfg":"FFFFFF","bbg":"E8F5E9","bfg":"1B2E1B","acc":"4CAF50","name":"🌿 Yashil"},
    "orange":  {"tbg":"BF360C","tfg":"FFFFFF","bbg":"FBE9E7","bfg":"2D1B0E","acc":"FF6D00","name":"🔥 Turunj"},
    "purple":  {"tbg":"4A148C","tfg":"FFFFFF","bbg":"EDE7F6","bfg":"1A0533","acc":"7B1FA2","name":"💜 Binafsha"},
    "minimal": {"tbg":"212121","tfg":"FFFFFF","bbg":"FFFFFF","bfg":"212121","acc":"BDBDBD","name":"⬜ Minimal"},
}
FREE_THEMES = {"classic", "green", "dark"}

# ════════════════════════════════════════════════
# DATABASE
# ════════════════════════════════════════════════
DB = "users_db.json"

def db_load():
    if not Path(DB).exists():
        return {}
    with open(DB, "r", encoding="utf-8") as f:
        return json.load(f)

def db_save(data):
    with open(DB, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

def today():
    return datetime.date.today().isoformat()

def get_user(uid):
    d = db_load(); k = str(uid)
    defaults = {
        "uid": uid, "username": "", "name": "", "premium": False,
        "day_count": 0, "total": 0, "last_date": "", "joined": today(),
        "refs": 0, "blocked": False
    }
    if k not in d:
        d[k] = defaults; db_save(d)
    else:
        changed = False
        migrations = [
            ("is_premium",   "premium"),
            ("daily_count",  "day_count"),
            ("total_count",  "total"),
            ("full_name",    "name"),
            ("referrals",    "refs"),
        ]
        for old, new in migrations:
            if old in d[k] and new not in d[k]:
                d[k][new] = d[k].pop(old); changed = True
            elif old in d[k]:
                d[k].pop(old); changed = True
        for key, val in defaults.items():
            if key not in d[k]:
                d[k][key] = val; changed = True
        if changed:
            db_save(d)
    return d[k]

def save_user(uid, **kw):
    d = db_load(); k = str(uid)
    if k not in d:
        get_user(uid); d = db_load()
    d[k].update(kw); db_save(d)

def reset_daily(uid):
    u = get_user(uid)
    if u["last_date"] != today():
        save_user(uid, day_count=0, last_date=today())

def can_use(uid):
    reset_daily(uid); u = get_user(uid)
    if u["blocked"]:
        return False, "blocked"
    lim = PREM_DAY if u["premium"] else FREE_DAY
    if u["day_count"] >= lim:
        return False, "limit"
    return True, "ok"

def add_usage(uid):
    u = get_user(uid)
    save_user(uid, day_count=u["day_count"]+1, total=u["total"]+1, last_date=today())

def all_users():
    return list(db_load().values())

# ════════════════════════════════════════════════
# YORDAMCHI FUNKSIYALAR
# ════════════════════════════════════════════════
def rgb(h):
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:], 16))

def is_admin(uid):
    return uid in ADMIN_IDS

def main_kb():
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("🆕  Yangi Prezentatsiya",   callback_data="new")],
        [InlineKeyboardButton("📄  Fayldan Prezentatsiya", callback_data="from_file")],
        [InlineKeyboardButton("💎  Premium",   callback_data="premium"),
         InlineKeyboardButton("📊  Statistika", callback_data="stats")],
        [InlineKeyboardButton("📖  Yordam",     callback_data="help")],
    ])

def back_kb():
    return InlineKeyboardMarkup([[
        InlineKeyboardButton("◀️  Bosh menyu", callback_data="back")
    ]])

# ════════════════════════════════════════════════
# FAYL MATNINI O'QISH FUNKSIYALARI
# ════════════════════════════════════════════════
def extract_text_from_pdf(file_bytes: bytes) -> tuple[str, str]:
    """PDF fayldan matn chiqaradi. Qaytaradi: (matn, xato_xabari)"""
    if not PYMUPDF_OK:
        return "", "PyMuPDF kutubxonasi o'rnatilmagan. `pip install pymupdf`"
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        texts = []
        total_chars = 0
        for page_num, page in enumerate(doc, 1):
            page_text = page.get_text()
            texts.append(page_text)
            total_chars += len(page_text.strip())
        doc.close()
        
        full_text = "\n".join(texts)
        
        # Agar juda kam matn bo'lsa - skanerlangan bo'lishi mumkin
        if total_chars < 100:
            return "", "PDF ichida deyarli matn yo'q - bu skanerlangan rasm bo'lishi mumkin (OCR kerak)"
        
        return full_text, ""
    except Exception as e:
        log.error("PDF o'qish xato: %s", e, exc_info=True)
        return "", f"PDF ochishda xato: {str(e)}"

def extract_text_from_pptx(file_bytes: bytes) -> tuple[str, str]:
    """PPTX fayldan matn chiqaradi. Qaytaradi: (matn, xato_xabari)"""
    try:
        prs = Presentation(BytesIO(file_bytes))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            # Markdown va format belgilarini tozalash
                            line = line.replace("**", "")  # Bold
                            line = line.replace("__", "")  # Underline
                            line = line.replace("~~", "")  # Strikethrough
                            line = re.sub(r'\*\*([^*]+)\*\*', r'\1', line)  # **bold**
                            line = re.sub(r'__([^_]+)__', r'\1', line)      # __italic__
                            line = re.sub(r'\*([^*]+)\*', r'\1', line)      # *italic*
                            line = re.sub(r'_([^_]+)_', r'\1', line)        # _italic_
                            texts.append(line)
        
        full_text = "\n".join(texts)
        if len(full_text.strip()) < 50:
            return "", "PPTX ichida deyarli matn yo'q"
        
        return full_text, ""
    except Exception as e:
        log.error("PPTX o'qish xato: %s", e, exc_info=True)
        return "", f"PPTX ochishda xato: {str(e)}"

def extract_text_from_docx(file_bytes: bytes) -> tuple[str, str]:
    """DOCX fayldan matn chiqaradi. Qaytaradi: (matn, xato_xabari)"""
    if not DOCX_OK:
        return "", "python-docx kutubxonasi o'rnatilmagan. `pip install python-docx`"
    try:
        doc = DocxDocument(BytesIO(file_bytes))
        texts = []
        for para in doc.paragraphs:
            line = para.text.strip()
            if line:
                # Markdown va format belgilarini tozalash
                line = line.replace("**", "")
                line = line.replace("__", "")
                line = line.replace("~~", "")
                line = re.sub(r'\*\*([^*]+)\*\*', r'\1', line)
                line = re.sub(r'__([^_]+)__', r'\1', line)
                line = re.sub(r'\*([^*]+)\*', r'\1', line)
                line = re.sub(r'_([^_]+)_', r'\1', line)
                texts.append(line)
        
        full_text = "\n".join(texts)
        if len(full_text.strip()) < 50:
            return "", "DOCX ichida deyarli matn yo'q"
        
        return full_text, ""
    except Exception as e:
        log.error("DOCX o'qish xato: %s", e, exc_info=True)
        return "", f"DOCX ochishda xato: {str(e)}"

def extract_text_from_file(file_bytes: bytes, mime_type: str, file_name: str) -> tuple[str, str, str]:
    """
    Fayl turini aniqlaydi va matn chiqaradi.
    Qaytaradi: (matn, fayl_turi, xato_xabari)
    """
    fname = file_name.lower()

    if fname.endswith(".pdf") or "pdf" in mime_type:
        text, error = extract_text_from_pdf(file_bytes)
        return text, "PDF", error

    elif fname.endswith(".pptx") or "presentationml" in mime_type or fname.endswith(".ppt"):
        text, error = extract_text_from_pptx(file_bytes)
        return text, "PPTX", error

    elif fname.endswith(".docx") or "wordprocessingml" in mime_type or fname.endswith(".doc"):
        text, error = extract_text_from_docx(file_bytes)
        return text, "DOCX", error

    return "", "UNKNOWN", "Noma'lum fayl turi"

# ════════════════════════════════════════════════
# PPTX YARATISH FUNKSIYALARI
# ════════════════════════════════════════════════
def pexels_img(q):
    if not PEXELS_KEY:
        return None
    try:
        r = requests.get(
            "https://api.pexels.com/v1/search",
            headers={"Authorization": PEXELS_KEY},
            params={"query": q, "per_page": 1, "orientation": "landscape"},
            timeout=8
        )
        photos = r.json().get("photos", [])
        if photos:
            return requests.get(photos[0]["src"]["medium"], timeout=8).content
    except Exception as e:
        log.warning("Pexels: %s", e)
    return None

def add_footer(slide, t, num, total):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(7.25), Inches(10), Inches(0.25))
    bar.fill.solid(); bar.fill.fore_color.rgb = rgb(t["acc"]); bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(8.8), Inches(7.26), Inches(1.1), Inches(0.2))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = "%d/%d" % (num, total)
    r.font.size = Pt(9); r.font.color.rgb = RGBColor(255, 255, 255); r.font.name = "Calibri"

def set_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = rgb(color)

def add_bullet(p, bullet, t, total_count):
    p.alignment = PP_ALIGN.LEFT
    spacing = Pt(2) if total_count >= 7 else Pt(5)
    p.space_before = spacing; p.space_after = spacing
    run = p.add_run()
    
    # Matn uzunligini cheklash (muvozanatli)
    MAX_BULLET_LEN = 180
    
    if bullet.strip().startswith("##"):
        text = bullet.strip().lstrip("#").strip()
        if len(text) > MAX_BULLET_LEN:
            text = text[:MAX_BULLET_LEN-3] + "..."
        run.text = text
        run.font.size  = Pt(13) if total_count >= 8 else Pt(15)
        run.font.bold  = True
        run.font.color.rgb = rgb(t["acc"])
    else:
        clean = bullet.strip().lstrip("-*>. ").strip()
        if len(clean) > MAX_BULLET_LEN:
            clean = clean[:MAX_BULLET_LEN-3] + "..."
        run.text = "  " + clean
        run.font.size  = Pt(11) if total_count >= 8 else Pt(12)
        run.font.bold  = False
        run.font.color.rgb = rgb(t["bfg"])
    run.font.name = "Calibri"

def make_pptx(topic, slides_data, theme_key="classic", add_img=False, lang="uz", source_label=""):
    t = THEMES.get(theme_key, THEMES["classic"])
    total = len(slides_data) + 2
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ─── Sarlavha slayd ───
    sl = prs.slides.add_slide(blank); set_bg(sl, t["tbg"])
    tx = sl.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(9), Inches(1.8))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = topic.upper()
    r.font.bold = True; r.font.size = Pt(32); r.font.name = "Calibri"
    r.font.color.rgb = rgb(t["tfg"])

    sub = sl.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(0.5))
    p2 = sub.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = ("📄 " + source_label + "  |  ") if source_label else ""
    r2.text += "AkademikYordamchi Bot  |  Gemini AI"
    r2.font.size = Pt(12); r2.font.color.rgb = rgb(t["acc"]); r2.font.name = "Calibri"

    dbox = sl.shapes.add_textbox(Inches(0.5), Inches(4.65), Inches(9), Inches(0.4))
    p3 = dbox.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = datetime.date.today().strftime("%d.%m.%Y")
    r3.font.size = Pt(12); r3.font.color.rgb = rgb(t["tfg"]); r3.font.name = "Calibri"

    acc = sl.shapes.add_shape(1, Inches(3.5), Inches(5.5), Inches(3), Pt(5))
    acc.fill.solid(); acc.fill.fore_color.rgb = rgb(t["acc"]); acc.line.fill.background()
    add_footer(sl, t, 1, total)

    # ─── Kontent slaydlar ───
    for i, sd in enumerate(slides_data, 2):
        sl = prs.slides.add_slide(blank); set_bg(sl, t["bbg"])
        bullets = sd.get("bullets", [])
        
        # Bullet sonini cheklash - lekin juda kam ham bo'lmasin
        total_chars = sum(len(b) for b in bullets)
        MAX_CHARS_PER_SLIDE = 1200  # Muvozanatli hajm
        
        if total_chars > MAX_CHARS_PER_SLIDE:
            # Bulletlarni qisqartirish - muhim bo'limlar va key pointlar
            filtered_bullets = []
            current_chars = 0
            for bullet in bullets:
                if current_chars + len(bullet) > MAX_CHARS_PER_SLIDE and len(filtered_bullets) >= 6:
                    break  # Kamida 6 ta bullet bo'lsin
                filtered_bullets.append(bullet)
                current_chars += len(bullet)
            bullets = filtered_bullets
            
            # Maksimum 10 ta bullet
            if len(bullets) > 10:
                bullets = bullets[:10]
        
        n_bullets = len(bullets)

        hpan = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.15))
        hpan.fill.solid(); hpan.fill.fore_color.rgb = rgb(t["tbg"]); hpan.line.fill.background()

        htx = sl.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(9.4), Inches(0.95))
        htf = htx.text_frame; htf.word_wrap = True
        hp = htf.paragraphs[0]; hp.alignment = PP_ALIGN.LEFT
        hr = hp.add_run(); hr.text = sd.get("title", "")
        hr.font.bold = True; hr.font.size = Pt(21); hr.font.name = "Calibri"
        hr.font.color.rgb = rgb(t["tfg"])

        has_img = False
        if add_img:
            img_bytes = pexels_img(sd.get("img_query", topic))
            if img_bytes and PIL_OK:
                try:
                    im = PILImage.open(BytesIO(img_bytes)).resize((310, 200))
                    buf = BytesIO(); im.save(buf, "JPEG"); buf.seek(0)
                    sl.shapes.add_picture(buf, Inches(6.6), Inches(1.25), Inches(3.1), Inches(2.0))
                    has_img = True
                except:
                    pass

        use_two_col = (n_bullets >= 7 and not has_img)
        if use_two_col:
            mid = (n_bullets + 1) // 2
            for col_l, col_w, col_bullets in [
                (Inches(0.35), Inches(4.55), bullets[:mid]),
                (Inches(5.1),  Inches(4.55), bullets[mid:])
            ]:
                btx = sl.shapes.add_textbox(col_l, Inches(1.25), col_w, Inches(5.85))
                tf = btx.text_frame; tf.word_wrap = True
                for j, b in enumerate(col_bullets):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    add_bullet(p, b, t, n_bullets)
        else:
            cw = Inches(6.0) if has_img else Inches(9.4)
            btx = sl.shapes.add_textbox(Inches(0.35), Inches(1.25), cw, Inches(5.85))
            tf = btx.text_frame; tf.word_wrap = True
            for j, b in enumerate(bullets):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                add_bullet(p, b, t, n_bullets)

        if sd.get("notes"):
            sl.notes_slide.notes_text_frame.text = sd["notes"]
        add_footer(sl, t, i, total)

    # ─── Yakuniy slayd ───
    sl = prs.slides.add_slide(blank); set_bg(sl, t["tbg"])
    tx = sl.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(9), Inches(1.4))
    p = tx.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    endings = {"uz": "E'tiboringiz uchun rahmat!", "ru": "Spasibo za vnimanie!", "en": "Thank you!"}
    r.text = endings.get(lang, endings["uz"])
    r.font.bold = True; r.font.size = Pt(30); r.font.name = "Calibri"
    r.font.color.rgb = rgb(t["tfg"])

    sub2 = sl.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
    p2 = sub2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "🤖 AkademikYordamchi Bot"
    r2.font.size = Pt(14); r2.font.color.rgb = rgb(t["acc"]); r2.font.name = "Calibri"
    add_footer(sl, t, total, total)

    safe = re.sub(r"[^\w\s-]", "", topic)[:28].strip().replace(" ", "_")
    fname = "prez_%s_%s.pptx" % (safe, datetime.datetime.now().strftime("%H%M%S"))
    prs.save(fname)
    return fname

# ════════════════════════════════════════════════
# AI KONTENT — MAVZUDAN
# ════════════════════════════════════════════════
async def ai_content(topic, n, lang="uz", style="university"):
    lang_map = {
        "uz": "O'zbek tilida yoz. Barcha matn faqat O'zbek tilida bo'lsin.",
        "ru": "Rus tilida yoz. Barcha matn faqat Rus tilida bo'lsin.",
        "en": "Write everything in English only.",
    }
    style_map = {
        "university": "Universitet darajasida ilmiy, professional, akademik uslubda. Atamalar, ta'riflar, misollar va tahlillar bilan boy qilib yoz.",
        "school":     "Maktab o'quvchilari uchun sodda, qiziqarli va tushunarli uslubda. Oddiy so'zlar, ko'p misol va solishtirish.",
        "business":   "Biznes prezentatsiyasi formatida. Statistika, foydalar, muammolar va yechimlar, raqamlar va natijalar bilan.",
        "creative":   "Ijodiy, qiziqarli va rang-barang uslubda. Qiziqarli faktlar, hikoyalar, misollar va o'xshatishlar bilan.",
    }
    prompt = """Sen dunyodagi eng yaxshi prezentatsiya mazmuni yozuvchi mutaxassissan.

VAZIFA: Quyidagi mavzu boyicha juda boy, batafsil va professional slaydlar tuzib ber.

Mavzu: {topic}
Til: {lang_str}
Uslub: {style_str}
Slayd soni: {n} ta

MUHIM QOIDALAR:
1. HAR BIR SLAYD uchun kamida 10-15 ta bullet yoz. Har bullet eng kamida 1 ta tolik, mazmuny gap bolsin.
2. ## belgisi bilan kichik bolim sarlavhasi qosh (har slaydda 2-3 ta).
3. Slaydlar bir-biridan butunlay farq qilsin va mavzuning turli tomonlarini ochsin.
4. Faktlar, raqamlar, misollar va izohlar bilan boyit.
5. img_query maydonini ingliz tilida yoz.
6. Aynan {n} ta slayd bolsin, na kam na ko'p.

FAQAT quyidagi JSON formatida javob qaytargin. Boshqa hech narsa yozma:
[
  {{
    "title": "Birinchi slayd sarlavhasi",
    "bullets": [
      "## Asosiy tushuncha va ta'rif",
      "Bu mavzuning to'liq va batafsil ta'rifi yoki izohi shu yerda yoziladi.",
      "Muhim jihat: aniq ma'lumot, raqam yoki fakt bilan mustahkamlangan gap.",
      "## Tarixiy rivojlanish yoki kelib chiqishi",
      "Bu bo'limdagi asosiy fikr to'liq yozilishi va misollar bilan boy qilinishi kerak.",
      "Qo'shimcha izoh yoki amaliy qo'llash misoli bu yerda joylashsin.",
      "Yana bir muhim fikr yoki xulosa."
    ],
    "notes": "Maruzachi uchun izoh: bu slaydda nimalarga e'tibor berish kerak.",
    "img_query": "relevant photo search query in english"
  }}
]""".format(
        topic=topic,
        lang_str=lang_map.get(lang, lang_map["uz"]),
        style_str=style_map.get(style, style_map["university"]),
        n=n
    )

    resp = ai.models.generate_content(model=MODEL, contents=[prompt])
    raw = resp.text.strip()
    raw = re.sub(r"^```[a-zA-Z]*\n?", "", raw)
    raw = re.sub(r"\n?```$", "", raw).strip()

    try:
        data = json.loads(raw)
        if isinstance(data, list) and len(data) > 0:
            return data[:n]
    except Exception as e:
        log.error("JSON parse xato: %s | raw[:300]: %s", e, raw[:300])

    return _fallback_slides(topic, n)

# ════════════════════════════════════════════════
# AI KONTENT — FAYLDAN  ← YaNGI FUNKSIYA
# ════════════════════════════════════════════════
async def ai_content_from_file(file_text: str, topic: str, n: int, lang: str = "uz", style: str = "university") -> list:
    """
    Yuklangan faylning matni asosida slayd kontenti yaratadi.
    file_text - fayldan chiqarilgan to'liq matn
    topic     - foydalanuvchi ko'rsatgan mavzu (fokus)
    """
    lang_map = {
        "uz": "O'zbek tilida yoz. Barcha matn faqat O'zbek tilida bo'lsin.",
        "ru": "Rus tilida yoz. Barcha matn faqat Rus tilida bo'lsin.",
        "en": "Write everything in English only.",
    }
    style_map = {
        "university": "Universitet darajasida ilmiy, professional, akademik uslubda. Atamalar, ta'riflar, misollar va tahlillar bilan boy qilib yoz.",
        "school":     "Maktab o'quvchilari uchun sodda, qiziqarli va tushunarli uslubda. Oddiy so'zlar, ko'p misol va solishtirish.",
        "business":   "Biznes prezentatsiyasi formatida. Statistika, foydalar, muammolar va yechimlar, raqamlar va natijalar bilan.",
        "creative":   "Ijodiy, qiziqarli va rang-barang uslubda. Qiziqarli faktlar, hikoyalar, misollar va o'xshatishlar bilan.",
    }

    # Matn juda uzun bo'lsa, qisqartirish (Gemini token limiti)
    MAX_CHARS = 18000
    trimmed_text = file_text[:MAX_CHARS]
    if len(file_text) > MAX_CHARS:
        trimmed_text += "\n\n[Matn uzunligi sababli qolgan qism qisqartirildi]"

    prompt = """Sen dunyodagi eng yaxshi prezentatsiya mazmuni yozuvchi mutaxassissan.

VAZIFA: Quyida berilgan HUJJAT MATNI asosida juda BOY, BATAFSIL va PROFESSIONAL prezentatsiya slaydlari tayyorla.

═══════════════════════════════════════════════════
HUJJAT MATNI:
═══════════════════════════════════════════════════
{file_text}
═══════════════════════════════════════════════════

Foydalanuvchi so'ragan MAVZU / FOKUS: {topic}
Til: {lang_str}
Uslub: {style_str}
Slayd soni: {n} ta

JUDA MUHIM QOIDALAR (QAT'IY BAJARING):

1. HUJJAT MATNINI TO'LIQ TAHLIL QILING:
   - Hujjatdagi BARCHA asosiy g'oyalarni, tushunchalarni va faktlarni qamrab oling
   - Har bir muhim bo'limni alohida slaydda yoritib bering
   - Hujjatdagi misollar, raqamlar, formulalar va ta'riflarni ALBATTA ishlatib bering

2. HAR BIR SLAYD BOY VA MAZMUNLI BO'LSIN:
   - 8-10 ta bullet point (kamida 8 ta bo'lishi kerak)
   - Har bir bullet 1-2 ta to'liq gap bo'lsin
   - Har bullet 120-160 belgi atrofida bo'lsin
   - Qisqa va sodda bulletlardan QOCHIB TURING

3. ## BO'LIM SARLAVHALARI:
   - Har slaydda 3 ta ## bo'lim sarlavhasi qo'shing
   - Har bo'limda 2-3 ta batafsil bullet bo'lsin
   - Masalan:
     ## Asosiy tushuncha va ta'rif
     Birinchi batafsil izoh yoki ta'rif (1-2 ta gap)
     Ikkinchi qo'shimcha ma'lumot yoki misol
     ## Amaliy qo'llash
     Amaliyotdagi foydalanish usullari batafsil
     Aniq misollar va holat tasvirlari
     ## Muhim xulosalar
     Asosiy natijalar va ta'kidlar

4. HUJJATDAGI MA'LUMOTLARNI BOYITING:
   - Hujjatdagi har bir muhim atamani to'liq tushuntiring
   - Agar hujjatda misol bo'lsa, uni batafsil yozib bering
   - Agar formula, algoritm yoki jarayon bo'lsa, bosqichma-bosqich tushuntiring
   - Raqamlar va statistika bo'lsa, ularni kontekst bilan bering

5. SLAYDLAR BIR-BIRIDAN FARQ QILSIN:
   - Har slayd hujjatning alohida jihatini yoritsin
   - Takrordan qoching - har slaydda yangi ma'lumot bo'lsin
   - Mantiqiy ketma-ketlik: kirish → asosiy tushunchalar → tahlil → xulosa

6. img_query maydonini ingliz tilida aniq va mavzuga mos yozing.

7. Aynan {n} ta slayd yarating - na kam, na ko'p.

8. FAQAT HUJJAT MATNIDAN FOYDANING - o'z boshingdan HECH NARSA QO'SHMANG!

JAVOB FORMATI - FAQAT JSON (boshqa hech narsa yozma):
[
  {{
    "title": "Birinchi slayd sarlavhasi (aniq va qisqa)",
    "bullets": [
      "## Birinchi bo'lim sarlavhasi",
      "Bu bo'limdagi birinchi batafsil izoh. To'liq 1-2 ta gap bo'lishi kerak.",
      "Ikkinchi muhim fikr yoki qo'shimcha tushuntirish. Yana bir gap.",
      "Uchinchi qo'shimcha ma'lumot yoki misol hujjat asosida.",
      "## Ikkinchi bo'lim sarlavhasi",
      "Bu bo'limdagi asosiy ma'lumot yoki konsepsiya batafsil yozilgan.",
      "Qo'shimcha kontekst yoki amaliy misol hujjat asosida berilgan.",
      "Muhim xulosalar yoki ta'kidlanishi kerak bo'lgan jihatlar.",
      "## Uchinchi bo'lim sarlavhasi",
      "Bu yerda yana alohida jihat yoki qo'shimcha ma'lumot yozilgan.",
      "Yakuniy fikr yoki xulosa shu mavzu bo'yicha."
    ],
    "notes": "Maruzachi uchun izoh: slaydda nimalarga e'tibor berish kerak.",
    "img_query": "specific and relevant photo search query in english"
  }}
]

ESLATMA: Hujjat katta bo'lgani uchun, barcha muhim qismlarini {n} ta slaydga TENG taqsimlang. Har slayd maksimal mazmunli va boy bo'lishi kerak!""".format(
        file_text=trimmed_text,
        topic=topic if topic else "umumiy tahlil",
        lang_str=lang_map.get(lang, lang_map["uz"]),
        style_str=style_map.get(style, style_map["university"]),
        n=n
    )

    resp = ai.models.generate_content(model=MODEL, contents=[prompt])
    raw = resp.text.strip()
    raw = re.sub(r"^```[a-zA-Z]*\n?", "", raw)
    raw = re.sub(r"\n?```$", "", raw).strip()

    try:
        data = json.loads(raw)
        if isinstance(data, list) and len(data) > 0:
            return data[:n]
    except Exception as e:
        log.error("Fayldan JSON parse xato: %s | raw[:300]: %s", e, raw[:300])

    return _fallback_slides(topic or "Hujjat", n)

def _fallback_slides(topic, n):
    """AI javob bermasa, fallback slaydlar."""
    return [
        {
            "title": "%s - %d-qism" % (topic, i + 1),
            "bullets": [
                "## Asosiy tushuncha",
                "%s mavzusining %d-qismi bo'yicha asosiy ma'lumot." % (topic, i + 1),
                "Bu mavzu juda muhim va keng qamrovli soha bo'lib, ko'plab jihatlarni o'z ichiga oladi.",
                "## Amaliy jihatlar",
                "Amaliyotda bu bilimlar turli sohalarda keng qo'llaniladi va muhim ahamiyatga ega.",
                "Mutaxassislar bu soha bo'yicha doimiy izlanishlar olib borishadi.",
                "Ko'proq ma'lumot olish uchun tegishli manbalar va adabiyotlarga murojaat qiling."
            ],
            "notes": "%s mavzusining %d-qismi." % (topic, i + 1),
            "img_query": topic
        }
        for i in range(n)
    ]

# ════════════════════════════════════════════════
# WIZARD — MAVZUDAN
# ════════════════════════════════════════════════
async def start_wizard(update: Update, ctx: ContextTypes.DEFAULT_TYPE, from_callback=False):
    uid = update.effective_user.id
    ok, reason = can_use(uid)
    if not ok:
        uu = get_user(uid)
        lim = PREM_DAY if uu["premium"] else FREE_DAY
        if reason == "blocked":
            msg = "🚫 *Siz bloklangansiz.*\n\nQo'shimcha ma'lumot uchun adminga murojaat qiling."
        else:
            msg = (
                "⚠️ *Bugungi limit tugadi!*\n\n"
                "📊 Ishlatildi: *%d / %d* ta\n\n"
                "━━━━━━━━━━━━━━━━━━━━━\n"
                "💎 *Premium* oling va:\n"
                "  • Kuniga *%d* ta yarating\n"
                "  • Cheklovsiz imkoniyatlar\n"
                "━━━━━━━━━━━━━━━━━━━━━\n\n"
                "🔗 /premium — batafsil ma'lumot\n"
                "⏰ _Ertaga avtomatik yangilanadi_"
            ) % (uu["day_count"], lim, PREM_DAY)
        if from_callback:
            await update.callback_query.edit_message_text(msg, parse_mode=ParseMode.MARKDOWN)
        else:
            await update.message.reply_text(msg, parse_mode=ParseMode.MARKDOWN)
        return

    ctx.user_data.clear()
    ctx.user_data["step"] = S_TOPIC
    ctx.user_data["mode"] = "topic"

    text = (
        "✏️ *Yangi Prezentatsiya*\n\n"
        "━━━━━━━━━━━━━━━━━━━━━\n"
        "📝 *Mavzuni yozing:*\n"
        "━━━━━━━━━━━━━━━━━━━━━\n\n"
        "💡 *Misollar:*\n"
        "  • Python dasturlash asoslari\n"
        "  • O'zbekiston tarixi\n"
        "  • Sun'iy intellekt va kelajak\n"
        "  • Marketing strategiyasi\n"
        "  • Iqlim o'zgarishi muammolari\n\n"
        "⌨️ _Mavzuni quyiga yozing..._"
    )
    if from_callback:
        await update.callback_query.edit_message_text(text, parse_mode=ParseMode.MARKDOWN)
    else:
        await update.message.reply_text(text, parse_mode=ParseMode.MARKDOWN)

# ════════════════════════════════════════════════
# WIZARD — FAYLDAN  ← YANGI
# ════════════════════════════════════════════════
async def start_file_wizard(update: Update, ctx: ContextTypes.DEFAULT_TYPE, from_callback=False):
    """Fayldan prezentatsiya yaratish wizardini boshlaydi."""
    uid = update.effective_user.id
    ok, reason = can_use(uid)
    if not ok:
        uu = get_user(uid)
        lim = PREM_DAY if uu["premium"] else FREE_DAY
        msg = (
            "⚠️ *Bugungi limit tugadi!*\n\n"
            "📊 Ishlatildi: *%d / %d* ta\n\n"
            "💎 /premium — Premium oling"
        ) % (uu["day_count"], lim) if reason == "limit" else "🚫 *Siz bloklangansiz.*"
        if from_callback:
            await update.callback_query.edit_message_text(msg, parse_mode=ParseMode.MARKDOWN)
        else:
            await update.message.reply_text(msg, parse_mode=ParseMode.MARKDOWN)
        return

    ctx.user_data.clear()
    ctx.user_data["step"] = "await_file"
    ctx.user_data["mode"] = "file"

    text = (
        "📄 *Fayldan Prezentatsiya*\n\n"
        "━━━━━━━━━━━━━━━━━━━━━\n"
        "📎 *Quyidagi fayl turlarini yuboring:*\n"
        "━━━━━━━━━━━━━━━━━━━━━\n\n"
        "  📕 *PDF* — darslik, maqola, hisobot\n"
        "  📊 *PPTX* — tayyor prezentatsiya\n"
        "  📝 *DOCX* — Word hujjati\n\n"
        "━━━━━━━━━━━━━━━━━━━━━\n"
        "🤖 Bot faylingizni tahlil qilib,\n"
        "    uning mazmuniga asoslanib\n"
        "    yangi slaydlar tayyorlaydi!\n"
        "━━━━━━━━━━━━━━━━━━━━━\n\n"
        "📎 _Faylni hozir yuboring..._"
    )
    if from_callback:
        await update.callback_query.edit_message_text(text, parse_mode=ParseMode.MARKDOWN)
    else:
        await update.message.reply_text(text, parse_mode=ParseMode.MARKDOWN)

# ════════════════════════════════════════════════
# FAYL HANDLER  ← YANGI
# ════════════════════════════════════════════════
async def on_file(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    """Foydalanuvchi PDF/PPTX/DOCX fayl yuborganda ishlaydi."""
    step = ctx.user_data.get("step", S_IDLE)

    # Faqat await_file rejimida fayl qabul qilinadi
    if step != "await_file":
        await update.message.reply_text(
            "📄 *Fayl qabul qilindi!*\n\n"
            "Agar fayldan prezentatsiya yaratmoqchi bo'lsangiz,\n"
            "avval 📄 *Fayldan Prezentatsiya* tugmasini bosing.",
            parse_mode=ParseMode.MARKDOWN,
            reply_markup=main_kb()
        )
        return

    doc = update.message.document
    if not doc:
        await update.message.reply_text(
            "⚠️ *Fayl topilmadi.*\nIltimos, PDF, PPTX yoki DOCX fayl yuboring.",
            parse_mode=ParseMode.MARKDOWN
        )
        return

    fname   = doc.file_name or "file"
    mime    = doc.mime_type or ""
    fsize   = doc.file_size or 0

    # Fayl hajmi tekshiruvi (max 20 MB)
    MAX_SIZE = 20 * 1024 * 1024
    if fsize > MAX_SIZE:
        await update.message.reply_text(
            "❌ *Fayl juda katta!*\n\n"
            "Maksimal hajm: *20 MB*\n"
            "Sizning faylingiz: *%.1f MB*" % (fsize / 1024 / 1024),
            parse_mode=ParseMode.MARKDOWN
        )
        return

    # Fayl turini tekshirish
    fname_lower = fname.lower()
    allowed_ext = (".pdf", ".pptx", ".ppt", ".docx", ".doc")
    if not any(fname_lower.endswith(e) for e in allowed_ext):
        await update.message.reply_text(
            "❌ *Noto'g'ri fayl turi!*\n\n"
            "Qabul qilinadigan formatlar:\n"
            "📕 PDF · 📊 PPTX · 📝 DOCX",
            parse_mode=ParseMode.MARKDOWN
        )
        return

    # Faylni yuklab olish
    status_msg = await update.message.reply_text(
        "⏳ *Fayl o'qilmoqda...*",
        parse_mode=ParseMode.MARKDOWN
    )

    try:
        await update.message.chat.send_action(ChatAction.TYPING)
        tg_file = await doc.get_file()
        file_bytes = await tg_file.download_as_bytearray()
        file_bytes = bytes(file_bytes)

        # Matnni chiqarish
        extracted_text, file_type, error_msg = extract_text_from_file(file_bytes, mime, fname)

        if error_msg or not extracted_text or len(extracted_text.strip()) < 50:
            error_detail = error_msg or "Fayl ichida matn topilmadi"
            await status_msg.edit_text(
                "❌ <b>Fayl matnini o'qib bo'lmadi!</b>\n\n"
                "📄 Fayl: <code>%s</code>\n"
                "📋 Tur: <b>%s</b>\n\n"
                "⚠️ <b>Sabab:</b>\n"
                "<i>%s</i>\n\n"
                "━━━━━━━━━━━━━━━━━━━━━\n"
                "<b>Yechimlar:</b>\n"
                "  • Agar PDF skanerlangan bo'lsa, OCR qilingan PDF yuboring\n"
                "  • Agar kutubxona yo'q bo'lsa, o'rnating\n"
                "  • Boshqa fayl yuboring\n"
                "  • /new bilan mavzu kiriting" % (fname[:40], file_type, error_detail),
                parse_mode=ParseMode.HTML
            )
            ctx.user_data.clear(); ctx.user_data["step"] = S_IDLE
            return

        # Ma'lumotlarni saqlash
        ctx.user_data["file_text"]  = extracted_text
        ctx.user_data["file_name"]  = fname
        ctx.user_data["file_type"]  = file_type
        ctx.user_data["step"]       = S_FILE_TOPIC

        word_count = len(extracted_text.split())
        await status_msg.edit_text(
            "✅ <b>Fayl muvaffaqiyatli o'qildi!</b>\n\n"
            "📄 Fayl: <code>%s</code>\n"
            "📋 Tur: <b>%s</b>\n"
            "📝 So'zlar: <b>~%d ta</b>\n\n"
            "━━━━━━━━━━━━━━━━━━━━━\n"
            "🎯 <b>Prezentatsiya mavzusini yozing:</b>\n"
            "━━━━━━━━━━━━━━━━━━━━━\n\n"
            "💡 Aniq mavzu ko'rsating yoki\n"
            "    «umumiy» deb yozing (barcha mazmun)\n\n"
            "⌨️ <i>Mavzuni yozing...</i>" % (fname[:30], file_type, word_count),
            parse_mode=ParseMode.HTML
        )

    except Exception as e:
        log.error("Fayl yuklash xato: %s", e, exc_info=True)
        await status_msg.edit_text(
            "❌ *Xato yuz berdi!*\n\n"
            "`%s`\n\n"
            "Qayta urinib ko'ring: /new" % str(e)[:200],
            parse_mode=ParseMode.MARKDOWN
        )
        ctx.user_data.clear(); ctx.user_data["step"] = S_IDLE

# ════════════════════════════════════════════════
# WIZARD — SLAYD SONI KB
# ════════════════════════════════════════════════
async def send_slides_kb(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    uu = get_user(uid); is_p = uu["premium"]
    topic = ctx.user_data.get("topic", "")
    mode  = ctx.user_data.get("mode", "topic")

    icon = "📄" if mode == "file" else "📌"
    kb = InlineKeyboardMarkup([
        [InlineKeyboardButton("📄  5 ta",  callback_data="w:slides:5"),
         InlineKeyboardButton("📄  8 ta",  callback_data="w:slides:8"),
         InlineKeyboardButton("📄 10 ta",  callback_data="w:slides:10")],
        [InlineKeyboardButton("📋 15 ta %s" % ("✅" if is_p else "💎"), callback_data="w:slides:15"),
         InlineKeyboardButton("📋 20 ta 💎",                             callback_data="w:slides:20")],
    ])
    await update.message.reply_text(
        "✅ *Qabul qilindi!*\n\n"
        "%s _%s_\n\n"
        "━━━━━━━━━━━━━━━━━━━━━\n"
        "📊 *Nechta slayd kerak?*\n"
        "━━━━━━━━━━━━━━━━━━━━━" % (icon, topic[:60]),
        parse_mode=ParseMode.MARKDOWN,
        reply_markup=kb
    )
    ctx.user_data["step"] = "slides"

async def send_style_kb(chat_id, ctx, bot, num):
    uu = get_user(int(chat_id)); is_p = uu["premium"]
    kb = InlineKeyboardMarkup([
        [InlineKeyboardButton("🎓 Universitet", callback_data="w:style:university"),
         InlineKeyboardButton("🏫 Maktab",      callback_data="w:style:school")],
        [InlineKeyboardButton("💼 Biznes %s" % ("✅" if is_p else "💎"), callback_data="w:style:business"),
         InlineKeyboardButton("🎨 Ijodiy %s"  % ("✅" if is_p else "💎"), callback_data="w:style:creative")],
    ])
    await bot.send_message(
        chat_id,
        "✅ *%d ta slayd* tanlandi!\n\n"
        "━━━━━━━━━━━━━━━━━━━━━\n"
        "🎨 *Qaysi uslubda bo'lsin?*\n"
        "━━━━━━━━━━━━━━━━━━━━━" % num,
        parse_mode=ParseMode.MARKDOWN,
        reply_markup=kb
    )
    ctx.user_data["step"] = "style"

async def send_theme_kb(chat_id, ctx, bot):
    uu = get_user(int(chat_id)); is_p = uu["premium"]
    rows = []
    items = list(THEMES.items())
    for i in range(0, len(items), 2):
        row = []
        for key, t in items[i:i+2]:
            label = t["name"] if (is_p or key in FREE_THEMES) else t["name"] + " 💎"
            row.append(InlineKeyboardButton(label, callback_data="w:theme:%s" % key))
        rows.append(row)
    await bot.send_message(
        chat_id,
        "━━━━━━━━━━━━━━━━━━━━━\n"
        "🖌️ *Dizayn sxemasini tanlang:*\n"
        "━━━━━━━━━━━━━━━━━━━━━",
        parse_mode=ParseMode.MARKDOWN,
        reply_markup=InlineKeyboardMarkup(rows)
    )
    ctx.user_data["step"] = "theme"

async def send_lang_kb(chat_id, ctx, bot):
    uu = get_user(int(chat_id)); is_p = uu["premium"]
    kb = InlineKeyboardMarkup([[
        InlineKeyboardButton("🇺🇿 O'zbek",                               callback_data="w:lang:uz"),
        InlineKeyboardButton("🇷🇺 Rus %s"    % ("✅" if is_p else "💎"),  callback_data="w:lang:ru"),
        InlineKeyboardButton("🇬🇧 Ingliz %s" % ("✅" if is_p else "💎"),  callback_data="w:lang:en"),
    ]])
    await bot.send_message(
        chat_id,
        "━━━━━━━━━━━━━━━━━━━━━\n"
        "🌐 *Til tanlang:*\n"
        "━━━━━━━━━━━━━━━━━━━━━",
        parse_mode=ParseMode.MARKDOWN,
        reply_markup=kb
    )
    ctx.user_data["step"] = "lang"

async def send_img_kb(chat_id, ctx, bot):
    uu = get_user(int(chat_id))
    if uu["premium"] and PEXELS_KEY:
        kb = InlineKeyboardMarkup([[
            InlineKeyboardButton("🖼️  Ha, rasm qo'sh", callback_data="w:img:yes"),
            InlineKeyboardButton("📄  Rasmsiz",          callback_data="w:img:no"),
        ]])
        await bot.send_message(
            chat_id,
            "━━━━━━━━━━━━━━━━━━━━━\n"
            "🖼️ *Slaydlarga rasm qo'shilsinmi?*\n"
            "_(Pexels.com dan avtomatik olinadi)_\n"
            "━━━━━━━━━━━━━━━━━━━━━",
            parse_mode=ParseMode.MARKDOWN,
            reply_markup=kb
        )
        ctx.user_data["step"] = "img"
    else:
        ctx.user_data["add_img"] = False
        ctx.user_data["step"] = S_IDLE
        await do_generate(chat_id, ctx, bot)

# ════════════════════════════════════════════════
# GENERATSIYA  ← YANGILANGAN (file rejimini qo'llab-quvvatlaydi)
# ════════════════════════════════════════════════
async def do_generate(chat_id, ctx, bot):
    uid = int(chat_id)
    ok, reason = can_use(uid)
    if not ok:
        await bot.send_message(chat_id, "⚠️ Limit tugadi. /premium")
        return

    topic     = ctx.user_data.get("topic", "Mavzu")
    n         = ctx.user_data.get("num_slides", 8)
    style     = ctx.user_data.get("style", "university")
    theme     = ctx.user_data.get("theme", "classic")
    lang      = ctx.user_data.get("lang", "uz")
    add_img   = ctx.user_data.get("add_img", False)
    mode      = ctx.user_data.get("mode", "topic")
    file_text = ctx.user_data.get("file_text", "")
    file_name = ctx.user_data.get("file_name", "")
    file_type = ctx.user_data.get("file_type", "")

    lang_flags = {"uz": "🇺🇿 O'zbek", "ru": "🇷🇺 Rus", "en": "🇬🇧 Ingliz"}
    source_info = ""
    if mode == "file" and file_name:
        source_info = "📄 %s (%s)" % (file_name[:25], file_type)

    await bot.send_chat_action(chat_id=chat_id, action=ChatAction.TYPING)
    status = await bot.send_message(
        chat_id,
        "⚙️ *Tayyorlanmoqda...*\n\n"
        "📌 _%s_\n"
        "%s"
        "📊 *%d ta* slayd\n"
        "🎨 *%s*\n"
        "🌐 *%s*\n"
        "🖼️ Rasm: *%s*\n\n"
        "⏳ _AI kontent yozmoqda..._" % (
            topic, ("📄 Manba: *%s*\n" % source_info) if source_info else "",
            n, THEMES[theme]["name"],
            lang_flags.get(lang, lang.upper()),
            "✅ Ha" if add_img else "❌ Yo'q"
        ),
        parse_mode=ParseMode.MARKDOWN
    )

    try:
        # Fayl yoki mavzu rejimiga qarab AI ni chaqirish
        if mode == "file" and file_text:
            slides = await ai_content_from_file(file_text, topic, n, lang, style)
        else:
            slides = await ai_content(topic, n, lang, style)

        await bot.edit_message_text(
            "🔧 *PPTX fayl yaratilmoqda...*\n\n"
            "📊 _%d ta slayd tayyorlanmoqda_" % len(slides),
            chat_id=chat_id, message_id=status.message_id,
            parse_mode=ParseMode.MARKDOWN
        )
        fname = make_pptx(topic, slides, theme, add_img, lang, source_info)
        add_usage(uid)
        uu = get_user(uid); lim = PREM_DAY if uu["premium"] else FREE_DAY
        remain = lim - uu["day_count"]

        # Caption yaratish - HTML format (Markdown muammosi yo'q)
        caption_text = (
            "╔══════════════════════════╗\n"
            "║  ✅  Prezentatsiya Tayyor ║\n"
            "╚══════════════════════════╝\n\n"
            "📌 <b>{topic}</b>\n"
            "{source_line}"
            "┌─────────────────────────\n"
            "│ 📊 Slaydlar: <b>{slides} ta</b>\n"
            "│ 🎨 Dizayn: <b>{design}</b>\n"
            "│ 🌐 Til: <b>{lang}</b>\n"
            "│ ⏳ Bugun qoldi: <b>{remain} ta</b>\n"
            "└─────────────────────────\n\n"
            "🔄 Yangi yaratish: /new\n"
            "📄 Fayldan: /fromfile"
        ).format(
            topic=topic,
            source_line=("📄 Manba: <i>{}</i>\n".format(source_info) if source_info else ""),
            slides=len(slides) + 2,
            design=THEMES[theme]["name"],
            lang=lang_flags.get(lang, lang.upper()),
            remain=remain
        )
        
        with open(fname, "rb") as f:
            await bot.send_document(
                chat_id, f,
                filename="%s.pptx" % topic[:40],
                caption=caption_text,
                parse_mode=ParseMode.HTML
            )
        os.remove(fname)
        await bot.delete_message(chat_id=chat_id, message_id=status.message_id)

    except Exception as e:
        log.error("Generatsiya xato: %s", e, exc_info=True)
        try:
            await bot.edit_message_text(
                "❌ *Xato yuz berdi!*\n\n"
                "`%s`\n\n"
                "🔄 Qayta urinib ko'ring: /new" % str(e)[:200],
                chat_id=chat_id, message_id=status.message_id,
                parse_mode=ParseMode.MARKDOWN
            )
        except:
            pass

    ctx.user_data.clear(); ctx.user_data["step"] = S_IDLE

# ════════════════════════════════════════════════
# COMMAND HANDLERS
# ════════════════════════════════════════════════
async def cmd_start(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    get_user(user.id)
    save_user(user.id, username=user.username or "", name=user.full_name or "")
    ctx.user_data.clear(); ctx.user_data["step"] = S_IDLE

    if ctx.args and ctx.args[0].isdigit():
        ref = int(ctx.args[0])
        if ref != user.id:
            ru = get_user(ref); save_user(ref, refs=ru["refs"] + 1)

    reset_daily(user.id); uu = get_user(user.id)
    lim = PREM_DAY if uu["premium"] else FREE_DAY

    await update.message.reply_text(
        "╔══════════════════════════════╗\n"
        "║   🎓  AkademikYordamchi     ║\n"
        "╚══════════════════════════════╝\n\n"
        "Assalomu alaykum, *%s*! 👋\n\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "🤖 *Nima qila olaman?*\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "📊  Professional prezentatsiya\n"
        "📄  PDF/PPTX/DOCX → Slayd\n"
        "🎨  6 xil chiroyli dizayn\n"
        "🌐  O'zbek · Rus · Ingliz tili\n"
        "📝  Boy AI mazmun (10-12 bullet)\n"
        "🗂️  2 ustunli layout\n"
        "📌  Ma'ruzachi izohlari (notes)\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n\n"
        "💳 Holat: *%s*\n"
        "📅 Bugun: *%d / %d* ta\n\n"
        "👇 *Boshlash uchun tugmani bosing:*" % (
            user.first_name,
            "💎 Premium" if uu["premium"] else "🆓 Bepul",
            uu["day_count"], lim
        ),
        parse_mode=ParseMode.MARKDOWN,
        reply_markup=main_kb()
    )

async def cmd_help(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "╔══════════════════════════════╗\n"
        "║      📖  Qo'llanma          ║\n"
        "╚══════════════════════════════╝\n\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "📌 *Asosiy buyruqlar:*\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "🏠 /start — Bosh menyu\n"
        "🆕 /new — Prezentatsiya (mavzudan)\n"
        "📄 /fromfile — Prezentatsiya (fayldan)\n"
        "💎 /premium — Premium ma'lumot\n"
        "📊 /stats — Mening hisobim\n"
        "❌ /cancel — Bekor qilish\n\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "🔢 *Mavzudan qanday ishlaydi?*\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "1️⃣  /new — mavzu kiriting\n"
        "2️⃣  Slayd sonini tanlang\n"
        "3️⃣  Uslub tanlang\n"
        "4️⃣  Dizayn tanlang\n"
        "5️⃣  Tilni tanlang\n"
        "6️⃣  PPTX faylni yuklab oling!\n\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "📄 *Fayldan qanday ishlaydi?*\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "1️⃣  /fromfile — buyrug'ini bosing\n"
        "2️⃣  PDF / PPTX / DOCX yuboring\n"
        "3️⃣  Mavzu/fokus yozing\n"
        "4️⃣  Slayd soni, uslub, dizayn, til\n"
        "5️⃣  PPTX faylni yuklab oling!\n\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "👑 *Admin buyruqlari:*\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "`/admin` · `/give_premium` · `/remove_premium`\n"
        "`/block` · `/unblock` · `/broadcast` · `/topusers`",
        parse_mode=ParseMode.MARKDOWN,
        reply_markup=main_kb()
    )

async def cmd_cancel(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    ctx.user_data.clear(); ctx.user_data["step"] = S_IDLE
    await update.message.reply_text(
        "❌ *Bekor qilindi.*\n\n"
        "Yangi prezentatsiya uchun /new bosing.",
        parse_mode=ParseMode.MARKDOWN,
        reply_markup=main_kb()
    )

async def cmd_stats(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id; reset_daily(uid); uu = get_user(uid)
    lim = PREM_DAY if uu["premium"] else FREE_DAY
    bot_info = await ctx.bot.get_me()
    ref_link = "https://t.me/%s?start=%d" % (bot_info.username, uid)
    await update.message.reply_text(
        "╔══════════════════════════════╗\n"
        "║    📊  Sizning Hisobingiz   ║\n"
        "╚══════════════════════════════╝\n\n"
        "👤 Ism: *%s*\n"
        "🆔 ID: `%d`\n"
        "📅 Ro'yxat: %s\n\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "💳 Holat: *%s*\n"
        "📈 Bugun: *%d / %d* ta\n"
        "🏆 Jami: *%d* ta\n"
        "👥 Referallar: *%d* ta\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n\n"
        "🔗 *Referal havolangiz:*\n"
        "`%s`" % (
            uu.get("name", "—"), uid, uu.get("joined", "—"),
            "💎 Premium" if uu["premium"] else "🆓 Bepul",
            uu["day_count"], lim, uu["total"], uu["refs"], ref_link
        ),
        parse_mode=ParseMode.MARKDOWN
    )

async def cmd_premium(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    bot_info = await ctx.bot.get_me()
    ref_link = "https://t.me/%s?start=%d" % (bot_info.username, uid)
    await update.message.reply_text(
        "╔══════════════════════════════╗\n"
        "║     💎  PREMIUM OBUNA       ║\n"
        "╚══════════════════════════════╝\n\n"
        "💰 Narx: *%s*\n\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "✨ *Premium imkoniyatlar:*\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "📊  Kuniga *%d* ta prezentatsiya\n"
        "📑  Maksimal *%d* ta slayd\n"
        "📄  PDF · PPTX · DOCX dan yaratish\n"
        "🖼️  Rasmlar (Pexels.com)\n"
        "🎨  6 xil professional dizayn\n"
        "🌐  3 ta til (UZ · RU · EN)\n"
        "💼  Biznes va ijodiy uslub\n"
        "📝  Ma'ruzachi izohlari\n\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "🆓 *Bepul versiya:*\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "📊  Kuniga *%d* ta\n"
        "📑  Maksimal *%d* ta slayd\n"
        "🎨  3 ta dizayn · faqat O'zbek\n\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "📲 *To'lash uchun:*\n"
        "`@YOUR_ADMIN_USERNAME`\n"
        "💳 Click · Payme · Uzum Bank\n\n"
        "👥 *Referal dasturi:* 🎁\n"
        "3 do'st = 1 oy Premium bepul!\n"
        "`%s`" % (
            PRICE, PREM_DAY, PREM_SLIDE, FREE_DAY, FREE_SLIDE, ref_link
        ),
        parse_mode=ParseMode.MARKDOWN
    )

# ════════════════════════════════════════════════
# CALLBACK HANDLER
# ════════════════════════════════════════════════
async def on_callback(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()
    data = q.data; uid = q.from_user.id; chat_id = q.message.chat_id

    if data == "new":
        await start_wizard(update, ctx, from_callback=True); return

    if data == "from_file":
        await start_file_wizard(update, ctx, from_callback=True); return

    if data == "premium":
        bot_info = await ctx.bot.get_me()
        ref_link = "https://t.me/%s?start=%d" % (bot_info.username, uid)
        await q.edit_message_text(
            "╔══════════════════════════════╗\n"
            "║     💎  PREMIUM OBUNA       ║\n"
            "╚══════════════════════════════╝\n\n"
            "💰 Narx: *%s*\n\n"
            "✨ *Premium:* %d ta/kun · max %d slayd\n"
            "🆓 *Bepul:* %d ta/kun · max %d slayd\n\n"
            "📄 Fayldan slayd · 🖼️ Rasmlar\n"
            "🎨 6 dizayn · 🌐 3 til\n\n"
            "📲 `@YOUR_ADMIN_USERNAME`\n\n"
            "👥 Referal: `%s`" % (
                PRICE, PREM_DAY, PREM_SLIDE, FREE_DAY, FREE_SLIDE, ref_link
            ),
            parse_mode=ParseMode.MARKDOWN,
            reply_markup=back_kb()
        ); return

    if data == "stats":
        reset_daily(uid); uu = get_user(uid); lim = PREM_DAY if uu["premium"] else FREE_DAY
        bot_info = await ctx.bot.get_me()
        ref_link = "https://t.me/%s?start=%d" % (bot_info.username, uid)
        await q.edit_message_text(
            "📊 *Sizning hisobingiz*\n\n"
            "💳 *%s*\n"
            "📈 Bugun: *%d / %d* ta\n"
            "🏆 Jami: *%d* ta\n"
            "👥 Referallar: *%d* ta\n\n"
            "🔗 `%s`" % (
                "💎 Premium" if uu["premium"] else "🆓 Bepul",
                uu["day_count"], lim, uu["total"], uu["refs"], ref_link
            ),
            parse_mode=ParseMode.MARKDOWN,
            reply_markup=back_kb()
        ); return

    if data == "help":
        await q.edit_message_text(
            "📖 *Qo'llanma*\n\n"
            "🆕 *Mavzudan:*\n"
            "1️⃣ /new → mavzu → slayd soni → uslub → dizayn → til → PPTX\n\n"
            "📄 *Fayldan:*\n"
            "1️⃣ /fromfile → fayl yuboring → mavzu → slayd soni → uslub → dizayn → til → PPTX\n\n"
            "💎 /premium · 📊 /stats",
            parse_mode=ParseMode.MARKDOWN,
            reply_markup=back_kb()
        ); return

    if data == "back":
        reset_daily(uid); uu = get_user(uid); lim = PREM_DAY if uu["premium"] else FREE_DAY
        await q.edit_message_text(
            "🏠 *Bosh Menyu*\n\n"
            "💳 *%s* | 📅 *%d / %d* ta" % (
                "💎 Premium" if uu["premium"] else "🆓 Bepul",
                uu["day_count"], lim
            ),
            parse_mode=ParseMode.MARKDOWN,
            reply_markup=main_kb()
        ); return

    # Wizard tugmalari
    if not data.startswith("w:"): return
    parts = data.split(":")
    if len(parts) < 3: return
    step = parts[1]; val = parts[2]
    uu = get_user(uid)

    if step == "slides":
        num = int(val)
        if num > FREE_SLIDE and not uu["premium"]:
            await q.answer("💎 %d ta slayd faqat Premium uchun!" % num, show_alert=True)
            return
        ctx.user_data["num_slides"] = num
        await send_style_kb(chat_id, ctx, ctx.bot, num)

    elif step == "style":
        if val in ("business", "creative") and not uu["premium"]:
            await q.answer("💎 Bu uslub faqat Premium uchun!", show_alert=True)
            return
        ctx.user_data["style"] = val
        await send_theme_kb(chat_id, ctx, ctx.bot)

    elif step == "theme":
        if val not in FREE_THEMES and not uu["premium"]:
            await q.answer("💎 Bu dizayn faqat Premium uchun!", show_alert=True)
            return
        ctx.user_data["theme"] = val
        await send_lang_kb(chat_id, ctx, ctx.bot)

    elif step == "lang":
        if val in ("ru", "en") and not uu["premium"]:
            await q.answer("💎 Bu til faqat Premium uchun!", show_alert=True)
            return
        ctx.user_data["lang"] = val
        await send_img_kb(chat_id, ctx, ctx.bot)

    elif step == "img":
        ctx.user_data["add_img"] = (val == "yes")
        ctx.user_data["step"] = S_IDLE
        await do_generate(chat_id, ctx, ctx.bot)

# ════════════════════════════════════════════════
# MATN HANDLER  ← YANGILANGAN
# ════════════════════════════════════════════════
async def on_message(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    step = ctx.user_data.get("step", S_IDLE)

    # 1. Mavzu kiritish (oddiy rejim)
    if step == S_TOPIC:
        topic = update.message.text.strip()
        if len(topic) < 3:
            await update.message.reply_text(
                "⚠️ *Juda qisqa!*\nKamida 3 ta harf kiriting:",
                parse_mode=ParseMode.MARKDOWN
            ); return
        if len(topic) > 200:
            await update.message.reply_text(
                "⚠️ *Juda uzun!*\n200 belgigacha kiriting:",
                parse_mode=ParseMode.MARKDOWN
            ); return
        ctx.user_data["topic"] = topic
        ctx.user_data["mode"]  = "topic"
        await send_slides_kb(update, ctx)

    # 2. Fayl yuklangandan keyin mavzu kiritish
    elif step == S_FILE_TOPIC:
        topic = update.message.text.strip()
        if topic.lower() in ("umumiy", "hammasi", "barchasi", "all", "общий"):
            topic = ""  # Umumiy — barcha mazmun
        if len(topic) > 200:
            await update.message.reply_text(
                "⚠️ *Juda uzun!*\n200 belgigacha kiriting:",
                parse_mode=ParseMode.MARKDOWN
            ); return
        # Mavzu bo'sh bo'lsa, fayl nomidan foydalanish
        display_topic = topic if topic else ctx.user_data.get("file_name", "Hujjat")
        ctx.user_data["topic"] = display_topic
        ctx.user_data["mode"]  = "file"
        await send_slides_kb(update, ctx)

    else:
        await update.message.reply_text(
            "💡 *Nima qilmoqchisiz?*\n\n"
            "🆕 Mavzudan → /new\n"
            "📄 Fayldan  → /fromfile",
            parse_mode=ParseMode.MARKDOWN,
            reply_markup=main_kb()
        )

# ════════════════════════════════════════════════
# ADMIN HANDLERS
# ════════════════════════════════════════════════
def admin_only(fn):
    @wraps(fn)
    async def wrap(update: Update, ctx, *a, **kw):
        if not is_admin(update.effective_user.id):
            await update.message.reply_text("🚫 *Admin emassiz!*", parse_mode=ParseMode.MARKDOWN); return
        return await fn(update, ctx, *a, **kw)
    return wrap

@admin_only
async def cmd_admin(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    users = all_users(); today_str = today()
    premium = sum(1 for x in users if x.get("premium"))
    blocked = sum(1 for x in users if x.get("blocked"))
    active  = sum(1 for x in users if x.get("last_date") == today_str and x.get("day_count", 0) > 0)
    total_p = sum(x.get("total", 0) for x in users)
    await update.message.reply_text(
        "╔══════════════════════════════╗\n"
        "║      👑  ADMIN PANEL        ║\n"
        "╚══════════════════════════════╝\n\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "👥 Foydalanuvchilar: *%d*\n"
        "💎 Premium: *%d*\n"
        "🚫 Bloklangan: *%d*\n"
        "📅 Bugun faol: *%d*\n"
        "📊 Jami yaratilgan: *%d* ta\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━\n\n"
        "📌 *Buyruqlar:*\n"
        "`/give_premium <id>` — Premium berish\n"
        "`/remove_premium <id>` — Premium olish\n"
        "`/block <id>` — Bloklash\n"
        "`/unblock <id>` — Blokdan chiqarish\n"
        "`/broadcast <matn>` — Xabar yuborish\n"
        "`/topusers` — Top 10 foydalanuvchi" % (
            len(users), premium, blocked, active, total_p
        ),
        parse_mode=ParseMode.MARKDOWN
    )

@admin_only
async def cmd_give(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if not ctx.args:
        await update.message.reply_text("📌 Foydalanish: `/give_premium <id>`", parse_mode=ParseMode.MARKDOWN); return
    try:
        tid = int(ctx.args[0]); save_user(tid, premium=True)
        await update.message.reply_text("✅ *%d* ga 💎 *Premium* berildi!" % tid, parse_mode=ParseMode.MARKDOWN)
        try:
            await ctx.bot.send_message(
                tid,
                "🎉 *Tabriklaymiz!*\n\n"
                "Sizga *Premium* obuna faollashtirildi! 💎\n\n"
                "🚀 /new yoki /fromfile bilan boshlang.",
                parse_mode=ParseMode.MARKDOWN
            )
        except:
            pass
    except:
        await update.message.reply_text("❌ *Noto'g'ri ID!*", parse_mode=ParseMode.MARKDOWN)

@admin_only
async def cmd_remove(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if not ctx.args:
        await update.message.reply_text("📌 Foydalanish: `/remove_premium <id>`", parse_mode=ParseMode.MARKDOWN); return
    try:
        tid = int(ctx.args[0]); save_user(tid, premium=False)
        await update.message.reply_text("✅ *%d* dan Premium olindi." % tid, parse_mode=ParseMode.MARKDOWN)
    except:
        await update.message.reply_text("❌ *Noto'g'ri ID!*", parse_mode=ParseMode.MARKDOWN)

@admin_only
async def cmd_block(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if not ctx.args:
        await update.message.reply_text("📌 Foydalanish: `/block <id>`", parse_mode=ParseMode.MARKDOWN); return
    try:
        tid = int(ctx.args[0]); save_user(tid, blocked=True)
        await update.message.reply_text("🚫 *%d* bloklandi." % tid, parse_mode=ParseMode.MARKDOWN)
    except:
        await update.message.reply_text("❌ *Noto'g'ri ID!*", parse_mode=ParseMode.MARKDOWN)

@admin_only
async def cmd_unblock(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if not ctx.args:
        await update.message.reply_text("📌 Foydalanish: `/unblock <id>`", parse_mode=ParseMode.MARKDOWN); return
    try:
        tid = int(ctx.args[0]); save_user(tid, blocked=False)
        await update.message.reply_text("✅ *%d* blokdan chiqarildi." % tid, parse_mode=ParseMode.MARKDOWN)
    except:
        await update.message.reply_text("❌ *Noto'g'ri ID!*", parse_mode=ParseMode.MARKDOWN)

@admin_only
async def cmd_broadcast(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if not ctx.args:
        await update.message.reply_text("📌 Foydalanish: `/broadcast <matn>`", parse_mode=ParseMode.MARKDOWN); return
    text = " ".join(ctx.args); users = all_users(); sent = failed = 0
    msg = await update.message.reply_text(
        "📢 *%d* ta foydalanuvchiga yuborilmoqda..." % len(users),
        parse_mode=ParseMode.MARKDOWN
    )
    for uu in users:
        if uu.get("blocked"): continue
        try:
            await ctx.bot.send_message(
                uu["uid"],
                "📢 *Xabar:*\n\n%s" % text,
                parse_mode=ParseMode.MARKDOWN
            )
            sent += 1; await asyncio.sleep(0.05)
        except:
            failed += 1
    await msg.edit_text(
        "✅ Yuborildi: *%d* ta\n❌ Xato: *%d* ta" % (sent, failed),
        parse_mode=ParseMode.MARKDOWN
    )

@admin_only
async def cmd_topusers(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    users = sorted(all_users(), key=lambda x: x.get("total", 0), reverse=True)[:10]
    medals = ["🥇","🥈","🥉","4️⃣","5️⃣","6️⃣","7️⃣","8️⃣","9️⃣","🔟"]
    lines = [
        "╔══════════════════════════════╗\n"
        "║       🏆  Top 10            ║\n"
        "╚══════════════════════════════╝\n"
    ]
    for i, uu in enumerate(users):
        nm = uu.get("name") or uu.get("username") or str(uu.get("uid", ""))
        lines.append("%s %s *%s* — %d ta" % (
            medals[i],
            "💎" if uu.get("premium") else "🆓",
            nm[:18], uu.get("total", 0)
        ))
    await update.message.reply_text("\n".join(lines), parse_mode=ParseMode.MARKDOWN)

# ════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════
def main():
    app = ApplicationBuilder().token(TOKEN).build()

    # Asosiy buyruqlar
    app.add_handler(CommandHandler("start",          cmd_start))
    app.add_handler(CommandHandler("help",           cmd_help))
    app.add_handler(CommandHandler("cancel",         cmd_cancel))
    app.add_handler(CommandHandler("stats",          cmd_stats))
    app.add_handler(CommandHandler("premium",        cmd_premium))
    app.add_handler(CommandHandler("new",            lambda u, c: start_wizard(u, c, from_callback=False)))
    app.add_handler(CommandHandler("fromfile",       lambda u, c: start_file_wizard(u, c, from_callback=False)))

    # Admin buyruqlar
    app.add_handler(CommandHandler("admin",          cmd_admin))
    app.add_handler(CommandHandler("give_premium",   cmd_give))
    app.add_handler(CommandHandler("remove_premium", cmd_remove))
    app.add_handler(CommandHandler("block",          cmd_block))
    app.add_handler(CommandHandler("unblock",        cmd_unblock))
    app.add_handler(CommandHandler("broadcast",      cmd_broadcast))
    app.add_handler(CommandHandler("topusers",       cmd_topusers))

    # Callback va xabar handlerlari
    app.add_handler(CallbackQueryHandler(on_callback))

    # Fayl handleri (PDF, PPTX, DOCX) — matn handleriga qadar joylashtirilishi shart
    app.add_handler(MessageHandler(
        filters.Document.ALL & ~filters.COMMAND,
        on_file
    ))

    # Matn handleri
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, on_message))

    async def post_init(application):
        await application.bot.set_my_commands([
            BotCommand("start",    "🏠 Bosh menyu"),
            BotCommand("new",      "🆕 Mavzudan prezentatsiya"),
            BotCommand("fromfile", "📄 Fayldan prezentatsiya"),
            BotCommand("premium",  "💎 Premium obuna"),
            BotCommand("stats",    "📊 Mening hisobim"),
            BotCommand("help",     "📖 Yordam"),
            BotCommand("cancel",   "❌ Bekor qilish"),
        ])
        log.info("Bot ishga tushdi! (v3.0)")

    app.post_init = post_init

    print("=" * 52)
    print("  🎓  AkademikYordamchi Bot v3.0")
    print("=" * 52)
    print("  🤖  Model     :", MODEL)
    print("  👑  Adminlar  :", ADMIN_IDS or "belgilanmagan")
    print("  🖼️   Pexels    :", "ulangan ✅" if PEXELS_KEY else "ulanmagan ⚠️")
    print("  📄  PDF       :", "✅" if PYMUPDF_OK else "⚠️ (pip install pymupdf)")
    print("  📝  DOCX      :", "✅" if DOCX_OK else "⚠️ (pip install python-docx)")
    print("=" * 52)
    app.run_polling(drop_pending_updates=True)


if __name__ == "__main__":
    main()